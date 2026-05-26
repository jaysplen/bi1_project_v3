#!/usr/bin/env python3
"""
Build PRESENTATION.pptx and PRESENTATION.pdf from the project's existing
analytical artefacts (exports/*.png + exports/metrics_summary.json).

The deck has 10 widescreen (16:9) slides and follows the assignment:
  1. Title
  2. CRM context and DM toolbox
  3-4. W1 - Inventory optimisation (question/data + result)
  5-6. W2 - Repair duration (question/data + result)
  7-8. W3 - QA failure triage (question/data + result)
  9. Conclusion / management recommendations
 10. Bibliography and frameworks

Run from the repository root:

    python3 build_presentation.py

Requires python-pptx and (for the PDF step) headless LibreOffice (`soffice`).
"""
from __future__ import annotations

import json
import shutil
import subprocess
import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


ROOT = Path(__file__).resolve().parent
METRICS_PATH = ROOT / "exports" / "metrics_summary.json"
OUT_PPTX = ROOT / "PRESENTATION.pptx"
OUT_PDF = ROOT / "PRESENTATION.pdf"

W1_DIR = ROOT / "exports" / "w1"
W2_DIR = ROOT / "exports" / "w2"
W3_DIR = ROOT / "exports" / "w3"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2E, 0x5F, 0xA1)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
SLATE = RGBColor(0x47, 0x55, 0x69)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Reference catalogue (matches ACADEMIC_METHOD_JUSTIFICATION.md numbering)
# ---------------------------------------------------------------------------
@dataclass(frozen=True)
class Reference:
    n: int
    short: str
    full: str
    url: str


REFERENCES: list[Reference] = [
    Reference(
        1, "Shearer 2000",
        "Shearer, C. (2000). The CRISP-DM model: The new blueprint for data mining. "
        "Journal of Data Warehousing, 5(4), 13-22.",
        "https://www.taylorfrancis.com/chapters/mono/10.1201/b12040-8/"
        "process-model-data-mining%E2%80%94crisp-dm-daniel-putler-robert-krider",
    ),
    Reference(
        2, "CRISP-DM 1.0",
        "CRISP-DM Consortium (2000). CRISP-DM 1.0: Step-by-step data mining guide.",
        "https://www.kde.cs.uni-kassel.de/wp-content/uploads/lehre/ws2015-16/"
        "kdd/files/CRISPWP-0800.pdf",
    ),
    Reference(
        3, "Agrawal & Srikant 1994",
        "Agrawal, R., & Srikant, R. (1994). Fast algorithms for mining association "
        "rules in large databases. Proc. VLDB '94, 487-499.",
        "https://snap.stanford.edu/class/cs224w-readings/Agrawal94AssosiationRule.pdf",
    ),
    Reference(
        4, "Brin et al. 1997",
        "Brin, S., Motwani, R., Ullman, J. D., & Tsur, S. (1997). Dynamic itemset "
        "counting and implication rules for market basket data. Proc. ACM SIGMOD, 255-276.",
        "https://dl.acm.org/doi/10.1145/253260.253325",
    ),
    Reference(
        5, "Mundler 2019",
        "Mundler, N. (2019). Association rule mining and itemset-correlation based "
        "variants. arXiv:1907.09535.",
        "https://arxiv.org/abs/1907.09535",
    ),
    Reference(
        6, "Wang et al. 2018",
        "Wang, J., Pan, X., Wang, L., & Wei, W. (2018). Method of spare parts "
        "prediction models evaluation based on grey comprehensive correlation degree "
        "and association rules mining: A case study in aviation. Mathematical Problems "
        "in Engineering, 2018, 2643405.",
        "https://doi.org/10.1155/2018/2643405",
    ),
    Reference(
        7, "Didriksen et al. 2026",
        "Didriksen, S. K., Sigsgaard, K. W., Mortensen, N. H., & Jespersen, C. B. "
        "(2026). Assigning spare parts management decision-making strategies: A "
        "holistic portfolio classification methodology. Applied Sciences, 16(4), 1961.",
        "https://doi.org/10.3390/app16041961",
    ),
    Reference(
        8, "Breiman et al. 1984",
        "Breiman, L., Friedman, J. H., Olshen, R. A., & Stone, C. J. (1984). "
        "Classification and Regression Trees. Chapman & Hall.",
        "https://www.taylorfrancis.com/books/mono/10.1201/9781315139470/"
        "classification-regression-trees-leo-breiman-jerome-friedman-olshen-charles-stone",
    ),
    Reference(
        9, "Rudin 2019",
        "Rudin, C. (2019). Stop explaining black box machine learning models for "
        "high stakes decisions and use interpretable models instead. Nature Machine "
        "Intelligence, 1, 206-215.",
        "https://arxiv.org/abs/1811.10154",
    ),
    Reference(
        10, "Jang et al. 2021",
        "Jang, J., Nana, D., Hochschild, J., & de Lorenzo, J. V. H. (2021). "
        "Predicting breakdown risk based on historical maintenance data for Air "
        "Force ground vehicles. arXiv:2112.13922.",
        "https://arxiv.org/abs/2112.13922",
    ),
    Reference(
        11, "Ding et al. 2022",
        "Ding, Y., Gao, A., Ryden, T., et al. (2022). Acela: Predictable "
        "datacenter-level maintenance job scheduling. arXiv:2212.05155.",
        "https://arxiv.org/abs/2212.05155",
    ),
    Reference(
        12, "Zhang et al. 2023",
        "Zhang, Y., Gao, Z., Sun, J., & Liu, L. (2023). Machine-learning algorithms "
        "for process condition data-based inclusion prediction in continuous-casting "
        "process. Sensors, 23(15), 6719.",
        "https://doi.org/10.3390/s23156719",
    ),
    Reference(
        13, "Sonawani & Mukhopadhyay 2013",
        "Sonawani, S., & Mukhopadhyay, D. (2013). A decision tree approach to "
        "classify web services using quality parameters. arXiv:1311.6240.",
        "https://arxiv.org/abs/1311.6240",
    ),
    Reference(
        14, "Ngai et al. 2009",
        "Ngai, E. W. T., Xiu, L., & Chau, D. C. K. (2009). Application of data "
        "mining techniques in customer relationship management: A literature review "
        "and classification. Expert Systems with Applications, 36(2), 2592-2602.",
        "https://doi.org/10.1016/j.eswa.2008.02.021",
    ),
]
REF_BY_N = {r.n: r for r in REFERENCES}


FRAMEWORKS = [
    ("scikit-learn (Decision Trees)", "https://scikit-learn.org/stable/modules/tree.html"),
    ("mlxtend (Apriori, association rules)",
     "https://rasbt.github.io/mlxtend/user_guide/frequent_patterns/apriori/"),
    ("pandas (data manipulation)", "https://pandas.pydata.org/docs/"),
    ("Jupyter (CRISP-DM notebooks)", "https://jupyter.org/"),
    ("CRISP-DM 1.0 guide",
     "https://www.kde.cs.uni-kassel.de/wp-content/uploads/lehre/ws2015-16/kdd/files/CRISPWP-0800.pdf"),
]


# ---------------------------------------------------------------------------
# Drawing helpers
# ---------------------------------------------------------------------------
def _new_pres() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=NAVY,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def _add_bullets(slide, x, y, w, h, items, *, size=14, color=SLATE,
                 bullet_color=BLUE, line_spacing=1.15):
    """items: list of strings; each becomes its own paragraph with a small disc bullet."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        p.line_spacing = line_spacing
        bullet = p.add_run()
        bullet.text = "\u25CF  "
        bullet.font.name = FONT
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = bullet_color
        body = p.add_run()
        body.text = item
        body.font.name = FONT
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return tb


def _add_hyperlink_run(paragraph, text, url, *, size=11, color=BLUE, bold=False):
    r = paragraph.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.underline = True
    r.hyperlink.address = url
    return r


def _header_band(slide, eyebrow: str, title: str):
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(1.1), NAVY)
    _add_text(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.35),
              eyebrow.upper(), size=11, bold=True, color=WHITE)
    _add_text(slide, Inches(0.5), Inches(0.42), Inches(12), Inches(0.6),
              title, size=24, bold=True, color=WHITE)


def _footer(slide, page: int, total: int = 10, note: str = ""):
    y = Inches(7.05)
    _add_rect(slide, Emu(0), y, SLIDE_W, Inches(0.45), LIGHT)
    _add_text(slide, Inches(0.5), Inches(7.13), Inches(8), Inches(0.3),
              "VacuTech \u00B7 DLMDSEBA02 \u00B7 BI & Data Mining (Phase 2)"
              + (f"  \u2022  {note}" if note else ""),
              size=10, color=SLATE)
    _add_text(slide, Inches(11.8), Inches(7.13), Inches(1.0), Inches(0.3),
              f"{page} / {total}", size=10, color=SLATE,
              align=PP_ALIGN.RIGHT)


def _footnote_band(slide, items: list[tuple[int, str]]):
    """items: list of (ref_number, short_label). Renders citations with hyperlinks."""
    if not items:
        return
    y = Inches(6.55)
    tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    sup_map = str.maketrans("0123456789", "\u2070\u00B9\u00B2\u00B3\u2074\u2075\u2076\u2077\u2078\u2079")
    for i, (n, label) in enumerate(items):
        ref = REF_BY_N[n]
        if i > 0:
            sep = p.add_run()
            sep.text = "   "
            sep.font.size = Pt(10)
        sup = p.add_run()
        sup.text = str(n).translate(sup_map) + " "
        sup.font.name = FONT
        sup.font.size = Pt(11)
        sup.font.bold = True
        sup.font.color.rgb = BLUE
        _add_hyperlink_run(p, label, ref.url, size=10, color=SLATE)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def slide_1_title(prs: Presentation, metrics: dict) -> None:
    slide = _add_blank(prs)
    _add_rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, NAVY)
    _add_rect(slide, Emu(0), Inches(3.05), SLIDE_W, Inches(0.05), BLUE)
    _add_text(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.5),
              "VACUTECH \u00B7 PHASE 2 \u00B7 BUSINESS INTELLIGENCE & DATA MINING",
              size=14, bold=True, color=BLUE, font=FONT)
    _add_text(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.9),
              "From RMAs to Recommendations",
              size=56, bold=True, color=WHITE, font=FONT)
    _add_text(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(0.8),
              "Three CRM questions, three transparent data-mining models",
              size=22, color=LIGHT, font=FONT)

    _add_text(slide, Inches(0.9), Inches(3.4), Inches(11.5), Inches(0.5),
              "Course module: DLMDSEBA02  \u2022  Composite presentation",
              size=14, color=LIGHT)

    w1 = metrics.get("w1", {})
    w2 = metrics.get("w2", {})
    w3 = metrics.get("w3", {})

    kpis = [
        ("W1 \u00B7 Apriori",
         f"{w1.get('rules_count', '—')} rules from {w1.get('pm_cases', '—')} PM cases",
         f"Pareto: {w1.get('pareto_cov80_parts', '—')} SKUs cover \u224880%"),
        ("W2 \u00B7 Decision Tree",
         f"R\u00B2 = {w2.get('r2', 0):.2f} on {w2.get('n_test', '—')} test cases",
         f"MAE \u2248 {w2.get('mae', 0):.2f} days at intake"),
        ("W3 \u00B7 DT classifier",
         f"Threshold = {w3.get('threshold', '—')}, recall {w3.get('recall_failure', 0):.0%}",
         f"Accuracy {w3.get('accuracy', 0):.0%} on {w3.get('n_test', '—')} test cases"),
    ]
    card_w = Inches(3.9)
    card_h = Inches(1.8)
    gap = Inches(0.25)
    start_x = Inches(0.9)
    y = Inches(4.5)
    for i, (h, l1, l2) in enumerate(kpis):
        x = start_x + (card_w + gap) * i
        _add_rect(slide, x, y, card_w, card_h, RGBColor(0x29, 0x3C, 0x66))
        _add_rect(slide, x, y, Inches(0.1), card_h, BLUE)
        _add_text(slide, x + Inches(0.3), y + Inches(0.2), card_w - Inches(0.4),
                  Inches(0.5), h, size=16, bold=True, color=WHITE)
        _add_text(slide, x + Inches(0.3), y + Inches(0.75), card_w - Inches(0.4),
                  Inches(0.5), l1, size=13, color=LIGHT)
        _add_text(slide, x + Inches(0.3), y + Inches(1.2), card_w - Inches(0.4),
                  Inches(0.5), l2, size=12, color=RGBColor(0xB4, 0xC2, 0xDB))

    _add_text(slide, Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.4),
              "All KPIs are read live from exports/metrics_summary.json",
              size=10, color=RGBColor(0xB4, 0xC2, 0xDB))


def slide_2_crm(prs: Presentation) -> None:
    slide = _add_blank(prs)
    _header_band(slide, "Context", "CRM tasks and the data-mining toolbox")

    intro = (
        "Customer Relationship Management aims to identify, attract, retain and "
        "develop customers. Data mining supports each of these CRM tasks with a "
        "different family of algorithms\u00B9\u207E\u00B9\u2074. VacuTech's three W "
        "questions sit on the retention and development side: stock the right parts, "
        "make credible promises, and prevent rework before customers feel it."
    )
    _add_text(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.1),
              intro, size=14, color=SLATE)

    quadrants = [
        ("Customer Identification", "Segmentation, target selection",
         "Clustering, classification"),
        ("Customer Attraction", "Direct marketing, lead scoring",
         "Classification, association rules"),
        ("Customer Retention",
         "Service quality, churn, complaint triage",
         "Classification (DT), regression"),
        ("Customer Development",
         "Cross-/up-sell, lifetime value, bundling",
         "Association rules, regression"),
    ]
    grid_x = Inches(0.5)
    grid_y = Inches(2.55)
    cell_w = Inches(5.9)
    cell_h = Inches(1.55)
    gap = Inches(0.2)
    for i, (title, what, algo) in enumerate(quadrants):
        cx = grid_x + (cell_w + gap) * (i % 2)
        cy = grid_y + (cell_h + gap) * (i // 2)
        _add_rect(slide, cx, cy, cell_w, cell_h, LIGHT)
        _add_rect(slide, cx, cy, Inches(0.12), cell_h, BLUE)
        _add_text(slide, cx + Inches(0.3), cy + Inches(0.15),
                  cell_w - Inches(0.4), Inches(0.4),
                  title, size=15, bold=True, color=NAVY)
        _add_text(slide, cx + Inches(0.3), cy + Inches(0.55),
                  cell_w - Inches(0.4), Inches(0.4),
                  what, size=12, color=SLATE)
        _add_text(slide, cx + Inches(0.3), cy + Inches(0.95),
                  cell_w - Inches(0.4), Inches(0.4),
                  "DM family: " + algo, size=12, bold=True, color=BLUE)

    mapping = (
        "VacuTech mapping  \u2022  W1 \u2192 Customer Development (parts bundling)   "
        "\u2022  W2 \u2192 Customer Retention (credible ETAs)   "
        "\u2022  W3 \u2192 Customer Retention (QA-failure prevention)"
    )
    _add_text(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.5),
              mapping, size=13, bold=True, color=NAVY)

    _footnote_band(slide, [(14, "Ngai et al. 2009"), (1, "Shearer 2000")])
    _footer(slide, 2)


def _wp_question_slide(prs, page, eyebrow, title, question, data_lines,
                       method_title, method_lines, hyperparams, fw_links,
                       citations):
    slide = _add_blank(prs)
    _header_band(slide, eyebrow, title)

    _add_text(slide, Inches(0.5), Inches(1.25), Inches(6.3), Inches(0.4),
              "Business question", size=13, bold=True, color=BLUE)
    _add_text(slide, Inches(0.5), Inches(1.6), Inches(6.3), Inches(1.3),
              question, size=15, color=NAVY)

    _add_text(slide, Inches(0.5), Inches(3.1), Inches(6.3), Inches(0.4),
              "Data used", size=13, bold=True, color=BLUE)
    _add_bullets(slide, Inches(0.5), Inches(3.45), Inches(6.3), Inches(2.8),
                 data_lines, size=13)

    # Right column: method box
    x = Inches(7.1)
    _add_rect(slide, x, Inches(1.25), Inches(5.75), Inches(5.0), LIGHT)
    _add_rect(slide, x, Inches(1.25), Inches(0.12), Inches(5.0), GREEN)
    _add_text(slide, x + Inches(0.3), Inches(1.4), Inches(5.4), Inches(0.4),
              method_title, size=15, bold=True, color=NAVY)
    _add_bullets(slide, x + Inches(0.3), Inches(1.85), Inches(5.4), Inches(1.7),
                 method_lines, size=12, color=SLATE)

    _add_text(slide, x + Inches(0.3), Inches(3.55), Inches(5.4), Inches(0.35),
              "Fixed hyperparameters", size=12, bold=True, color=BLUE)
    _add_bullets(slide, x + Inches(0.3), Inches(3.9), Inches(5.4), Inches(1.4),
                 hyperparams, size=12, color=SLATE)

    _add_text(slide, x + Inches(0.3), Inches(5.3), Inches(5.4), Inches(0.35),
              "Frameworks", size=12, bold=True, color=BLUE)
    tb = slide.shapes.add_textbox(x + Inches(0.3), Inches(5.65), Inches(5.4),
                                  Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    for i, (label, url) in enumerate(fw_links):
        if i > 0:
            sep = p.add_run()
            sep.text = "   \u00B7   "
            sep.font.size = Pt(11)
            sep.font.color.rgb = SLATE
        _add_hyperlink_run(p, label, url, size=11)

    _footnote_band(slide, citations)
    _footer(slide, page)


def _wp_result_slide(prs, page, eyebrow, title, img_left, img_right,
                     kpi_rows, recommendations, citations):
    slide = _add_blank(prs)
    _header_band(slide, eyebrow, title)

    img_y = Inches(1.3)
    img_w = Inches(6.0)
    img_h = Inches(3.5)
    slide.shapes.add_picture(str(img_left), Inches(0.5), img_y,
                             width=img_w, height=img_h)
    slide.shapes.add_picture(str(img_right), Inches(6.85), img_y,
                             width=img_w, height=img_h)

    # KPI strip
    kpi_y = Inches(4.95)
    kpi_h = Inches(0.85)
    n = len(kpi_rows)
    total_w = Inches(12.33)
    cell_w = Emu(int(total_w / n))
    for i, (label, value) in enumerate(kpi_rows):
        cx = Inches(0.5) + Emu(int(cell_w) * i)
        _add_rect(slide, cx, kpi_y, cell_w - Inches(0.1), kpi_h, LIGHT)
        _add_text(slide, cx + Inches(0.2), kpi_y + Inches(0.08),
                  cell_w - Inches(0.4), Inches(0.32),
                  label, size=11, bold=True, color=BLUE)
        _add_text(slide, cx + Inches(0.2), kpi_y + Inches(0.38),
                  cell_w - Inches(0.4), Inches(0.45),
                  value, size=16, bold=True, color=NAVY)

    _add_text(slide, Inches(0.5), Inches(5.95), Inches(2.5), Inches(0.35),
              "Management action", size=13, bold=True, color=BLUE)
    _add_bullets(slide, Inches(3.0), Inches(5.95), Inches(9.8), Inches(0.65),
                 recommendations, size=12, color=SLATE, line_spacing=1.05)

    _footnote_band(slide, citations)
    _footer(slide, page)


def slide_3_w1_q(prs: Presentation) -> None:
    _wp_question_slide(
        prs, page=3,
        eyebrow="Work package 1 \u00B7 Inventory optimisation",
        title="Which add-on parts should be stocked or bundled for PM visits?",
        question=(
            "During preventive-maintenance jobs, technicians often consume parts "
            "that are NOT in the standard PM kit \u2014 forcing a second trip and "
            "eroding customer trust. We want a short, ranked stocking shortlist "
            "and a set of bundle rules a planner can act on."
        ),
        data_lines=[
            "538 PM cases from repairs.csv (maintenance_kit_applied = 1)",
            "Add-on lines from parts_used.csv (kit_part_flag = 0)",
            "Pivoted into a boolean basket matrix: 538 cases x 36 add-on parts",
            "No imputation, no scaling \u2014 Apriori works on boolean baskets",
        ],
        method_title="Apriori association rules + Pareto",
        method_lines=[
            "Apriori (mlxtend) discovers frequent co-occurring add-on parts",
            "Rules ranked by lift, retained when lift exceeds independence",
            "Pareto chart on part counts complements rules with a SKU shortlist",
        ],
        hyperparams=[
            "min_support = 0.05  (\u2248 25-50 baskets at this dataset size)",
            "min_lift = 1.2  (meaningfully stronger than chance)",
            "Outputs: w1_rules.csv, w1_metrics.json, two PNG charts",
        ],
        fw_links=[
            ("mlxtend.apriori",
             "https://rasbt.github.io/mlxtend/user_guide/frequent_patterns/apriori/"),
            ("mlxtend.association_rules",
             "https://rasbt.github.io/mlxtend/user_guide/frequent_patterns/association_rules/"),
        ],
        citations=[(3, "Agrawal & Srikant 1994"), (4, "Brin et al. 1997")],
    )


def slide_4_w1_r(prs: Presentation, metrics: dict) -> None:
    w1 = metrics.get("w1", {})
    _wp_result_slide(
        prs, page=4,
        eyebrow="Work package 1 \u00B7 Result",
        title="Bundle rules + a 21-SKU Pareto shortlist",
        img_left=W1_DIR / "w1_top_rules.png",
        img_right=W1_DIR / "w1_pareto.png",
        kpi_rows=[
            ("Rules found", str(w1.get("rules_count", "—"))),
            ("Top lift", f"{w1.get('top_lift', 0):.2f}"),
            ("Top confidence", f"{w1.get('top_confidence', 0):.0%}"),
            ("Pareto 80% / 90%",
             f"{w1.get('pareto_cov80_parts', '—')} / {w1.get('pareto_cov90_parts', '—')}"),
        ],
        recommendations=[
            "Lock the Pareto-head SKUs at every regional warehouse",
            "Pilot one bundle rule per pump series next to the standard PM kit",
            "Track the 'first-visit completion' rate before vs after roll-out",
        ],
        citations=[(6, "Wang et al. 2018"), (7, "Didriksen et al. 2026")],
    )


def slide_5_w2_q(prs: Presentation) -> None:
    _wp_question_slide(
        prs, page=5,
        eyebrow="Work package 2 \u00B7 Repair-duration prediction",
        title="At intake, how long will this repair take?",
        question=(
            "The service planner must commit to a customer-facing window. The "
            "model must use only intake-safe features (no QA or post-repair "
            "fields) and remain explainable to a front-desk advisor."
        ),
        data_lines=[
            "Target: repair_duration_days (right-skewed)",
            "Numeric: pump_age_years, technician_experience_years, parts_cost_eur",
            "Categorical: pump_model, complexity_class, failure_type, "
            "parts_from_hq_flag, region",
            "Preparation: dropna + one-hot encoding (no scaler)",
        ],
        method_title="Decision Tree Regressor (CART)",
        method_lines=[
            "One interpretable algorithm \u2014 splits readable on a whiteboard",
            "Trees are scale-free and need no target transform",
            "80/20 train/test split, seed = 7, evaluated on held-out data",
        ],
        hyperparams=[
            "max_depth = 6  (captures pump x complexity interactions)",
            "min_samples_leaf = 20  (\u2265 20 cases per leaf)",
            "random_state = 7  (deterministic split)",
        ],
        fw_links=[
            ("scikit-learn DecisionTreeRegressor",
             "https://scikit-learn.org/stable/modules/generated/sklearn.tree.DecisionTreeRegressor.html"),
            ("CART (Breiman et al. 1984)",
             "https://www.taylorfrancis.com/books/mono/10.1201/9781315139470/classification-regression-trees-leo-breiman-jerome-friedman-olshen-charles-stone"),
        ],
        citations=[(8, "Breiman et al. 1984"), (9, "Rudin 2019")],
    )


def slide_6_w2_r(prs: Presentation, metrics: dict) -> None:
    w2 = metrics.get("w2", {})
    _wp_result_slide(
        prs, page=6,
        eyebrow="Work package 2 \u00B7 Result",
        title="Usable intake-time ETAs \u2014 publish as a +/- range",
        img_left=W2_DIR / "w2_actual_vs_pred.png",
        img_right=W2_DIR / "w2_feature_importance.png",
        kpi_rows=[
            ("Train / Test", f"{w2.get('n_train', '—')} / {w2.get('n_test', '—')}"),
            ("R\u00B2", f"{w2.get('r2', 0):.3f}"),
            ("MAE (days)", f"{w2.get('mae', 0):.2f}"),
            ("RMSE (days)", f"{w2.get('rmse', 0):.2f}"),
        ],
        recommendations=[
            "Publish predicted duration as a \u00B12-day range on the customer acknowledgement",
            "Use feature importances as a routing signal (long jobs \u2192 senior techs)",
            "Refresh the model quarterly; track MAE drift",
        ],
        citations=[(10, "Jang et al. 2021"), (11, "Ding et al. 2022 \u00B7 Acela")],
    )


def slide_7_w3_q(prs: Presentation) -> None:
    _wp_question_slide(
        prs, page=7,
        eyebrow="Work package 3 \u00B7 QA-failure risk",
        title="Which finished repairs deserve a closer pre-QA look?",
        question=(
            "Before QA sign-off, operations wants a single yes/no flag per job. "
            "The target is imbalanced (~15% failures), so we use a deliberate "
            "operating threshold instead of bare accuracy."
        ),
        data_lines=[
            "Target: qa_failed_flag (binary; \u224815% positives)",
            "Numeric: pump_age_years, technician_experience_years, "
            "parts_cost_eur, repair_duration_days",
            "Categorical: pump_model, complexity_class, failure_type, "
            "parts_from_hq_flag",
            "Stratified 80/20 split on the failure flag",
        ],
        method_title="Decision Tree Classifier + fixed threshold",
        method_lines=[
            "predict_proba -> single operating threshold = 0.30",
            "One algorithm, one operating point = governable dashboard flag",
            "Evaluation: confusion matrix, recall + precision (not accuracy alone)",
        ],
        hyperparams=[
            "max_depth = 6 ; min_samples_leaf = 20 ; random_state = 7",
            "threshold = 0.30 (balances recall vs review capacity)",
            "stratify = y to preserve class balance in train/test",
        ],
        fw_links=[
            ("scikit-learn DecisionTreeClassifier",
             "https://scikit-learn.org/stable/modules/generated/sklearn.tree.DecisionTreeClassifier.html"),
            ("Imbalanced classification metrics",
             "https://scikit-learn.org/stable/modules/model_evaluation.html#classification-metrics"),
        ],
        citations=[(8, "Breiman et al. 1984"), (9, "Rudin 2019")],
    )


def slide_8_w3_r(prs: Presentation, metrics: dict) -> None:
    w3 = metrics.get("w3", {})
    _wp_result_slide(
        prs, page=8,
        eyebrow="Work package 3 \u00B7 Result",
        title="A governable pre-QA triage flag at threshold 0.30",
        img_left=W3_DIR / "w3_confusion_matrix.png",
        img_right=W3_DIR / "w3_feature_importance.png",
        kpi_rows=[
            ("Threshold", f"{w3.get('threshold', '—')}"),
            ("Recall (failure)", f"{w3.get('recall_failure', 0):.0%}"),
            ("Precision (failure)", f"{w3.get('precision_failure', 0):.0%}"),
            ("Accuracy", f"{w3.get('accuracy', 0):.0%}"),
        ],
        recommendations=[
            "Route flagged jobs through a 5-minute senior-tech pre-QA review",
            "Monitor the flag-vs-failure ratio weekly and adjust the threshold",
            "Avoid post-repair-only features when scoring at intake-time triage",
        ],
        citations=[(12, "Zhang et al. 2023"), (13, "Sonawani & Mukhopadhyay 2013")],
    )


def slide_9_conclusion(prs: Presentation, metrics: dict) -> None:
    slide = _add_blank(prs)
    _header_band(slide, "Conclusion",
                 "Three concrete recommendations for VacuTech management")

    w1 = metrics.get("w1", {})
    w2 = metrics.get("w2", {})
    w3 = metrics.get("w3", {})

    cards = [
        (
            "W1 \u00B7 Inventory",
            f"Pre-stock the top {w1.get('pareto_cov80_parts', '—')} add-on SKUs "
            f"at every regional warehouse and pilot one bundle rule per pump "
            f"series. Apriori found {w1.get('rules_count', '—')} rules with "
            f"top lift {w1.get('top_lift', 0):.1f}.",
            "CRM lens: customer development \u2014 first-visit completion + cross-sell of bundled BOMs.",
            BLUE,
        ),
        (
            "W2 \u00B7 Repair ETA",
            f"Publish duration on the customer acknowledgement as a "
            f"\u00B12-day range. The shallow tree reaches R\u00B2 "
            f"\u2248 {w2.get('r2', 0):.2f} and MAE \u2248 "
            f"{w2.get('mae', 0):.2f} days, while staying readable to the "
            f"front-desk advisor.",
            "CRM lens: customer retention \u2014 credible promises, fewer status calls.",
            GREEN,
        ),
        (
            "W3 \u00B7 QA triage",
            f"Send the jobs flagged at threshold {w3.get('threshold', '—')} "
            f"through a 5-minute senior pre-QA review. Recall "
            f"\u2248 {w3.get('recall_failure', 0):.0%} at "
            f"{w3.get('accuracy', 0):.0%} overall accuracy makes the "
            f"workload manageable.",
            "CRM lens: customer retention \u2014 prevent repeat-RMA situations.",
            RED,
        ),
    ]
    card_w = Inches(4.07)
    card_h = Inches(4.5)
    gap = Inches(0.2)
    start_x = Inches(0.5)
    y = Inches(1.4)
    for i, (h, body, crm, accent) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        _add_rect(slide, x, y, card_w, card_h, LIGHT)
        _add_rect(slide, x, y, card_w, Inches(0.08), accent)
        _add_text(slide, x + Inches(0.3), y + Inches(0.25),
                  card_w - Inches(0.5), Inches(0.5),
                  h, size=18, bold=True, color=NAVY)
        _add_text(slide, x + Inches(0.3), y + Inches(0.95),
                  card_w - Inches(0.5), Inches(2.6),
                  body, size=13, color=SLATE)
        _add_rect(slide, x + Inches(0.3), y + Inches(3.45),
                  card_w - Inches(0.6), Inches(0.04), accent)
        _add_text(slide, x + Inches(0.3), y + Inches(3.6),
                  card_w - Inches(0.5), Inches(0.85),
                  crm, size=12, bold=True, color=accent)

    why = (
        "Why interpretable models? Stakeholders (planners, advisors, QA "
        "supervisors) must act on the output. A shallow, auditable model that "
        "they can read and override is preferable to a high-accuracy black "
        "box they cannot challenge\u2079\u207B\u00B2."
    )
    _add_text(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.6),
              why, size=12, color=NAVY)

    _footnote_band(slide,
                   [(9, "Rudin 2019"), (2, "CRISP-DM Consortium 2000")])
    _footer(slide, 9)


def slide_10_bibliography(prs: Presentation) -> None:
    slide = _add_blank(prs)
    _header_band(slide, "Bibliography", "Peer-reviewed sources and frameworks")

    # Bibliography
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.25), Inches(8.3),
                                  Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    for i, ref in enumerate(REFERENCES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        p.line_spacing = 1.05
        n = p.add_run()
        n.text = f"[{ref.n}] "
        n.font.name = FONT
        n.font.size = Pt(9.5)
        n.font.bold = True
        n.font.color.rgb = BLUE
        body = p.add_run()
        body.text = ref.full + "  "
        body.font.name = FONT
        body.font.size = Pt(9.5)
        body.font.color.rgb = SLATE
        _add_hyperlink_run(p, "link", ref.url, size=9.5, color=BLUE)

    # Frameworks panel
    fx = Inches(9.0)
    _add_rect(slide, fx, Inches(1.25), Inches(3.83), Inches(5.5), LIGHT)
    _add_rect(slide, fx, Inches(1.25), Inches(0.12), Inches(5.5), GREEN)
    _add_text(slide, fx + Inches(0.3), Inches(1.4), Inches(3.4), Inches(0.4),
              "Frameworks used", size=15, bold=True, color=NAVY)

    tb2 = slide.shapes.add_textbox(fx + Inches(0.3), Inches(1.85),
                                   Inches(3.4), Inches(4.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Emu(0)
    tf2.margin_top = Emu(0)
    for i, (label, url) in enumerate(FRAMEWORKS):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(8)
        bullet = p.add_run()
        bullet.text = "\u25CF  "
        bullet.font.size = Pt(11)
        bullet.font.color.rgb = GREEN
        _add_hyperlink_run(p, label, url, size=11, color=NAVY, bold=True)

    _add_text(slide, fx + Inches(0.3), Inches(6.0), Inches(3.4), Inches(0.6),
              "See ACADEMIC_METHOD_JUSTIFICATION.md for full per-W "
              "argumentation.",
              size=9, color=SLATE)

    _footer(slide, 10)


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------
def build_pptx(metrics: dict, out: Path) -> Path:
    prs = _new_pres()
    slide_1_title(prs, metrics)
    slide_2_crm(prs)
    slide_3_w1_q(prs)
    slide_4_w1_r(prs, metrics)
    slide_5_w2_q(prs)
    slide_6_w2_r(prs, metrics)
    slide_7_w3_q(prs)
    slide_8_w3_r(prs, metrics)
    slide_9_conclusion(prs, metrics)
    slide_10_bibliography(prs)
    prs.save(out)
    return out


def build_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    if not shutil.which("soffice"):
        print("  [warn] LibreOffice (`soffice`) not found; skipping PDF step.")
        return False
    try:
        subprocess.run(
            [
                "soffice", "--headless",
                "--convert-to", "pdf",
                "--outdir", str(pdf_path.parent),
                str(pptx_path),
            ],
            check=True,
            capture_output=True,
            cwd=ROOT,
        )
    except subprocess.CalledProcessError as exc:
        print(f"  [warn] LibreOffice conversion failed: {exc.stderr.decode(errors='ignore')}")
        return False
    produced = pdf_path.parent / (pptx_path.stem + ".pdf")
    if produced.is_file() and produced != pdf_path:
        produced.replace(pdf_path)
    return pdf_path.is_file()


def main() -> int:
    if not METRICS_PATH.is_file():
        print(
            "ERROR: missing exports/metrics_summary.json \u2014 run "
            "`python3 bi_pipeline.py` first.",
            file=sys.stderr,
        )
        return 1
    metrics = json.loads(METRICS_PATH.read_text(encoding="utf-8"))

    pptx_path = build_pptx(metrics, OUT_PPTX)
    print(f"  PPTX: {pptx_path.name} ({pptx_path.stat().st_size // 1024} KB)")

    if build_pdf(pptx_path, OUT_PDF):
        print(f"  PDF:  {OUT_PDF.name} ({OUT_PDF.stat().st_size // 1024} KB)")
        return 0
    print(
        "ERROR: PDF build failed. Install LibreOffice (`soffice`) to enable "
        "the PPTX -> PDF step.",
        file=sys.stderr,
    )
    return 1


if __name__ == "__main__":
    sys.exit(main())
